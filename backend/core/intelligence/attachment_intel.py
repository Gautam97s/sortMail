"""
Attachment Intelligence
-----------------------
Analyzes document attachments (PDF, DOCX, PPTX, XLSX, images) and produces deep analysis.
Uses Gemini Flash to extract summary, key points, action items, financial data, etc.
"""

import json
import logging
import os
from typing import Optional, Dict, Any
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy.future import select

from app.config import settings
from models.attachment import Attachment, AttachmentStatus
import re
import asyncio
from core.intelligence.llama_engine import _call_llama
from core.intelligence.prompts import (
    ATTACHMENT_ANALYSIS_SYSTEM_PROMPT,
    ATTACHMENT_ANALYSIS_USER_PROMPT_TEMPLATE,
)
from core.intelligence.attachment_security import (
    precheck_attachment_metadata,
    sanitize_text_for_llm,
    scan_text_security,
    text_quality_stats,
)


logger = logging.getLogger(__name__)


async def analyze_attachment(attachment_id: str, db: AsyncSession) -> Optional[Dict[str, Any]]:
    """
    Extracts text from a downloaded attachment and runs a unified Gemini Flash analysis.
    Saves the result to attachment.intel_json and returns the JSON payload.
    """
    stmt = select(Attachment).where(Attachment.id == attachment_id)
    result = await db.execute(stmt)
    attachment = result.scalars().first()
    
    if not attachment:
        logger.warning(f"Attachment {attachment_id} not found.")
        return None
        
    if not attachment.storage_path:
        logger.warning(f"Attachment {attachment_id} has no storage path yet.")
        return None

    allowed, reason = precheck_attachment_metadata(
        filename=attachment.filename or attachment.filename_sanitized or "",
        mime_type=attachment.mime_type or "",
        size_bytes=int(attachment.size_bytes or 0),
        max_size_mb=settings.MAX_ATTACHMENT_SIZE_MB,
    )
    if not allowed:
        attachment.status = AttachmentStatus.QUARANTINED
        attachment.skip_reason = reason
        await db.commit()
        logger.warning(f"Attachment {attachment_id} quarantined: {reason}")
        return None

    if not os.path.exists(attachment.storage_path):
        attachment.status = AttachmentStatus.FAILED
        attachment.skip_reason = "file_missing_on_disk"
        await db.commit()
        logger.warning(f"Attachment {attachment_id} file missing on disk.")
        return None

    # 1. Extract Document Text
    text = await _extract_text(attachment)
    if not text or len(text.strip()) < 10:
        logger.debug(f"Attachment {attachment_id} text extraction failed or insufficient.")
        attachment.status = AttachmentStatus.FAILED
        await db.commit()
        return None
        
    # Keep the strongest signal: head + tail, capped to a compact window.
    if len(text) > 25000:
        safe_text = text[:20000] + "\n\n[...truncated...]\n\n" + text[-5000:]
    else:
        safe_text = text

    security_scan = scan_text_security(safe_text)
    safe_text = sanitize_text_for_llm(safe_text)

    quality = text_quality_stats(safe_text)
    if not quality["looks_like_text"]:
        attachment.status = AttachmentStatus.SKIPPED
        attachment.skip_reason = "insufficient_readable_text"
        attachment.extracted_text = safe_text[:1000]
        await db.commit()
        logger.info(f"Attachment {attachment_id} skipped due to low text quality.")
        return None

    # 2. Build Deep Analysis Prompt
    prompt = ATTACHMENT_ANALYSIS_USER_PROMPT_TEMPLATE.format(document_text=safe_text)

    messages = [
        {"role": "system", "content": ATTACHMENT_ANALYSIS_SYSTEM_PROMPT},
        {"role": "user", "content": prompt}
    ]
    
    for attempt in range(4):
        try:
            raw_json = await _call_llama(
                messages,
                max_tokens=1600,
                temperature=0.1,
                operation="attachment_intel",
                metadata={
                    "user_id": attachment.user_id,
                    "related_entity_type": "attachment",
                    "related_entity_id": attachment_id,
                },
            )
            
            # Clean up any potential markdown fences
            raw_json = re.sub(r'^```(?:json)?\s*|\s*```$', '', raw_json.strip(), flags=re.DOTALL).strip()
            # Extract JSON if there's surrounding text
            json_match = re.search(r'\{.*\}', raw_json, re.DOTALL)
            if json_match:
                raw_json = json_match.group(0)
            
            intel_dict = json.loads(raw_json)
            intel_dict["security"] = {
                **security_scan,
                "text_quality": quality,
            }
            
            # 3. Save back to DB
            attachment.intel_json = intel_dict
            attachment.extracted_text = safe_text[:2000] # Save a snippet
            attachment.status = AttachmentStatus.INDEXED
            
            await db.commit()
            logger.info(f"✅ Generated intelligence for attachment {attachment.id}")
            return intel_dict
            
        except json.JSONDecodeError as e:
            logger.warning(f"Attachment analysis JSON parse failed attempt {attempt+1}/4 for {attachment.id}: {e}")
            if attempt == 3:
                break
            await asyncio.sleep(1)
            
        except Exception as e:
            error_str = str(e)
            logger.warning(f"Attachment analysis failed attempt {attempt+1}/4 for {attachment.id}: {error_str}")
            
            if "429" in error_str or "rate" in error_str.lower():
                wait = 2 ** attempt
                logger.warning(f"Rate limited — waiting {wait}s")
                await asyncio.sleep(wait)
                continue
                
            if attempt == 3:
                break
            await asyncio.sleep(1)
            
    # Fallback / Failure
    logger.error(f"Failed to analyze attachment {attachment.id} after 4 attempts.")
    attachment.status = AttachmentStatus.FAILED
    await db.commit()
    return None


# ─────────────────────────────────────────────────────────────────────────────
# Text Extractors
# ─────────────────────────────────────────────────────────────────────────────

async def _extract_text(attachment: Attachment) -> str:
    """Route text extraction based on mime type or extension."""
    try:
        mime = (attachment.mime_type or "").lower()
        path = attachment.storage_path
        
        if mime == "application/pdf" or path.endswith(".pdf"):
            return _extract_pdf_text(path)
        elif "wordprocessing" in mime or path.endswith(".docx"):
            return _extract_docx_text(path)
        elif "presentation" in mime or path.endswith(".pptx"):
            return _extract_pptx_text(path)
        elif "spreadsheetml" in mime or path.endswith((".xlsx", ".xlsm")):
            return _extract_xlsx_text(path)
        elif mime.startswith("text/"):
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        elif mime.startswith("image/") or path.lower().endswith((".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff")):
            return _extract_image_ocr(path)
        else:
            logger.debug(f"Unsupported mime type for text extraction: {mime}")
            return ""
    except Exception as e:
        logger.error(f"Text extraction failed for {attachment.id}: {e}")
        return ""


def _extract_pdf_text(path: str) -> str:
    try:
        from pypdf import PdfReader
        reader = PdfReader(path)
        return "\\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
    except ImportError:
        logger.warning("pypdf not installed.")
        return ""
    except Exception as e:
        logger.warning(f"Failed parsing PDF {path}: {e}")
        return ""


def _extract_docx_text(path: str) -> str:
    try:
        from docx import Document
        doc = Document(path)
        return "\\n".join([p.text for p in doc.paragraphs if p.text])
    except ImportError:
        logger.warning("python-docx not installed.")
        return ""
    except Exception as e:
        logger.warning(f"Failed parsing DOCX {path}: {e}")
        return ""


def _extract_pptx_text(path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text.append(shape.text)
        return "\\n".join(text)
    except ImportError:
        logger.warning("python-pptx not installed.")
        return ""
    except Exception as e:
        logger.warning(f"Failed parsing PPTX {path}: {e}")
        return ""


def _extract_xlsx_text(path: str) -> str:
    try:
        from openpyxl import load_workbook

        wb = load_workbook(path, read_only=True, data_only=True)
        lines: list[str] = []

        for sheet in wb.worksheets:
            lines.append(f"[Sheet: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                cells = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                if cells:
                    lines.append(" | ".join(cells))

        wb.close()
        return "\n".join(lines)
    except ImportError:
        logger.warning("openpyxl not installed.")
        return ""
    except Exception as e:
        logger.warning(f"Failed parsing XLSX {path}: {e}")
        return ""


def _extract_image_ocr(path: str) -> str:
    try:
        from PIL import Image
        import pytesseract

        image = Image.open(path)
        text = pytesseract.image_to_string(image)
        return text or ""
    except ImportError:
        logger.debug("OCR dependencies (Pillow/pytesseract) not installed.")
        return ""
    except Exception as e:
        logger.warning(f"Failed OCR parse for image {path}: {e}")
        return ""
